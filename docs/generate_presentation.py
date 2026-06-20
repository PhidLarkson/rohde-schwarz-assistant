#!/usr/bin/env python3
"""Generate the hackathon presentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Brand colors
DARK = RGBColor(0x0A, 0x25, 0x40)
CYAN = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF6, 0xF8)
DIM = RGBColor(0x6B, 0x70, 0x85)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x00, 0x6C, 0xBE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(bg_color=DARK):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = bg_color
    return slide

def text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def bullet_box(slide, left, top, width, height, items, font_size=16, color=WHITE, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return tf

def add_bar(slide, left, top, width, height, color=CYAN):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 1.5, 11, 1.5, "RHODA", 72, CYAN, True)
text_box(s, 0.8, 3.0, 11, 1, "Virtual Lab Assistant for R&S Oscilloscopes", 32, WHITE)
add_bar(s, 0.8, 4.2, 3, 0.04, CYAN)
text_box(s, 0.8, 4.5, 11, 0.8, "Track 2: Virtual Lab Assistant (Conversational / Voice)", 20, CYAN)
text_box(s, 0.8, 5.5, 11, 1.2, "Rohde & Schwarz AI-Assisted Onboarding Hackathon\nKwame Nkrumah University of Science and Technology\nJune 2026", 16, DIM)
text_box(s, 0.8, 6.6, 11, 0.6, "Team: Afari  ·  Gregory  ·  William  ·  Prince", 14, DIM)

# ═══════════════════════════════════════════════
# SLIDE 2: Problem
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "THE PROBLEM", 14, CYAN, True)
text_box(s, 0.8, 1.0, 11, 1.5, "New students struggle with oscilloscope procedures", 40, WHITE, True)
bullet_box(s, 0.8, 2.8, 5.5, 4, [
    "→  Novices lack accessible step-by-step guidance",
    "→  Translating lab procedures into commands is error-prone",
    "→  Safety mistakes risk equipment damage and injury",
    "→  No feedback on what went wrong or why",
    "→  Instructors can't track individual student progress",
], 18, RGBColor(0xB0, 0xB8, 0xCC))
bullet_box(s, 7, 2.8, 5.5, 4, [
    "What students need:",
    "",
    "• A guide that speaks their language",
    "• Step-by-step walkthrough of procedures",
    "• Safety checks before touching hardware",
    "• Real-time instrument feedback",
    "• Adaptive learning that tracks progress",
], 18, WHITE)

# ═══════════════════════════════════════════════
# SLIDE 3: Solution
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "OUR SOLUTION", 14, CYAN, True)
text_box(s, 0.8, 1.0, 11, 1.2, "Rhoda — conversational AI lab assistant", 36, WHITE, True)
bullet_box(s, 0.8, 2.5, 5.8, 4.5, [
    "Voice + Text Chat Interface",
    "   Ask questions naturally, get step-by-step guidance",
    "",
    "Safety-Gated Instrument Control",
    "   Every WRITE command requires explicit confirmation",
    "",
    "AI-Powered Troubleshooting",
    "   Describe a bad trace → get diagnosis + fix steps",
    "",
    "XR Avatar (Quest Headset)",
    "   3D assistant with lip sync and hand tracking",
], 17, WHITE)
bullet_box(s, 7.2, 2.5, 5.5, 4.5, [
    "Real R&S RTB24 Integration",
    "   SCPI control via RsInstrument over USB",
    "",
    "Guided Workflows",
    "   2 demo flows: measurement + safety procedures",
    "",
    "RAG Knowledge Base",
    "   15 documents from RTB2 manual + procedures",
    "",
    "Progress Tracking",
    "   9 competency areas, curriculum-ordered learning",
], 17, WHITE)

# ═══════════════════════════════════════════════
# SLIDE 4: Architecture
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "SYSTEM ARCHITECTURE", 14, CYAN, True)

# Frontend box
add_bar(s, 0.5, 1.3, 6, 0.04, CYAN)
text_box(s, 0.5, 1.4, 6, 0.5, "FRONTEND (TypeScript + Vite)", 14, CYAN, True)
bullet_box(s, 0.7, 1.9, 5.5, 2.2, [
    "XR Avatar — Three.js + WebXR + IWSDK",
    "Text Chat — /chat.html lab console",
    "Instructor — /instructor.html dashboard",
    "Voice — Gemini STT/TTS + GhanaNLP (Twi)",
    "RAG — 18 embedded procedure/safety docs",
    "Anomaly — fault pattern matching + Gemini",
], 13, RGBColor(0xB0, 0xB8, 0xCC))

# Backend box
add_bar(s, 0.5, 4.3, 6, 0.04, CYAN)
text_box(s, 0.5, 4.4, 6, 0.5, "BACKEND (Flask + Python)", 14, CYAN, True)
bullet_box(s, 0.7, 4.9, 5.5, 2, [
    "13 REST API endpoints",
    "RAG retrieval (sentence-transformers)",
    "Instrument API with confirmation gate",
    "Session logging + progress tracking",
    "Failure point analytics",
], 13, RGBColor(0xB0, 0xB8, 0xCC))

# Right side — tech stack
add_bar(s, 7, 1.3, 5.5, 0.04, CYAN)
text_box(s, 7, 1.4, 5.5, 0.5, "TECHNOLOGY STACK", 14, CYAN, True)
bullet_box(s, 7.2, 1.9, 5, 5, [
    "Three.js + WebXR     3D rendering, XR",
    "Meta IWSDK           hand tracking, panels",
    "Google Gemini 2.5    LLM, vision, STT, TTS",
    "GhanaNLP API         Twi speech + translation",
    "Flask 3.0            REST backend",
    "sentence-transformers  RAG embeddings",
    "RsInstrument        R&S VISA/SCPI driver",
    "",
    "METRICS",
    "6,720 lines of code",
    "22 TypeScript + 6 Python modules",
    "15 RAG corpus documents",
    "45 GLB animation assets",
    "11 automated backend tests",
    "13 API endpoints",
    "3 web entry points",
], 13, WHITE)

# ═══════════════════════════════════════════════
# SLIDE 5: Demo Flow — Workflow 1
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "DEMO WORKFLOW 1: SIMPLE", 14, CYAN, True)
text_box(s, 0.8, 1.0, 11, 1, "Measure a 1 kHz Sine Wave (10 steps)", 32, WHITE, True)

steps_left = [
    "1.  Connect probe to CH1 BNC input",
    "2.  Set probe attenuation to 10X",
    "3.  Enable CH1",
    "4.  Set vertical scale to 1 V/div",
    "5.  Set timebase to 500 µs/div",
]
steps_right = [
    "6.  Set trigger: CH1, rising edge, auto",
    "7.  Adjust trigger level to midpoint",
    "8.  Verify waveform is stable",
    "9.  Read measurement (expect ~1 kHz)",
    "10. Review results",
]
bullet_box(s, 0.8, 2.5, 5.5, 4, steps_left, 17, WHITE)
bullet_box(s, 7, 2.5, 5.5, 4, steps_right, 17, WHITE)

text_box(s, 0.8, 5.8, 11, 1, "Each instrument-changing step triggers a confirmation prompt.\nThe student must approve before the oscilloscope setting changes.", 16, RGBColor(0xFF, 0xB5, 0x47))

# ═══════════════════════════════════════════════
# SLIDE 6: Demo Flow — Workflow 2
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, RED)
text_box(s, 0.8, 0.4, 10, 0.8, "DEMO WORKFLOW 2: SAFETY-CRITICAL", 14, RED, True)
text_box(s, 0.8, 1.0, 11, 1, "Handle an Overvoltage Warning (6 steps)", 32, WHITE, True)

bullet_box(s, 0.8, 2.5, 11, 4.5, [
    "⚠  Step 1:  STOP — Disconnect probe from the circuit immediately",
    "              Overvoltage can damage the instrument and create a safety hazard",
    "",
    "⚠  Step 2:  Check signal voltage with a multimeter before reconnecting",
    "",
    "⚠  Step 3:  Verify voltage is within safe limits (300V CAT II with 10X probe)",
    "",
    "    Step 4:  Select correct probe attenuation (10X or 100X for high voltage)",
    "",
    "    Step 5:  Set oscilloscope to maximum V/div range before reconnecting",
    "",
    "    Step 6:  Reconnect probe carefully, verify no overload warning",
], 17, WHITE)

# ═══════════════════════════════════════════════
# SLIDE 7: Safety & Confirmation Gate
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "SAFETY ARCHITECTURE", 14, CYAN, True)
text_box(s, 0.8, 1.0, 11, 1, "Every instrument change requires explicit approval", 32, WHITE, True)

bullet_box(s, 0.8, 2.5, 5.5, 4.5, [
    "Confirmation Gate",
    "",
    "1. User or Rhoda requests a change",
    "2. Backend generates preview with:",
    "   • Current value",
    "   • Proposed value",
    "   • Human-readable summary",
    "   • Unique confirmation token",
    "3. Student sees prompt and decides",
    "4. Only confirmed changes execute",
    "5. Token expires if state changes",
], 16, WHITE)

bullet_box(s, 7, 2.5, 5.5, 4.5, [
    "Safety Rules",
    "",
    "• READ operations never need confirmation",
    "• WRITE operations always need confirmation",
    "• Parameter safety limits enforced",
    "  (V/div: 1mV–10V, timebase: 1ns–50s)",
    "• Overvoltage detection → immediate warning",
    "• Default to simulator mode",
    "• Real hardware only with SCOPE_RESOURCE",
    "• No PII in session logs",
    "• Anomaly detection flags unsafe traces",
], 16, WHITE)

# ═══════════════════════════════════════════════
# SLIDE 8: RTB24 Integration
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "REAL HARDWARE: R&S RTB24", 14, CYAN, True)
text_box(s, 0.8, 1.0, 11, 1, "SCPI commands verified against RTB2 User Manual", 28, WHITE, True)

bullet_box(s, 0.8, 2.3, 5.8, 4.8, [
    "Connection (via RsInstrument Python library):",
    "  pip install RsInstrument pyvisa-py",
    "  SCOPE_RESOURCE=USB python server/app.py",
    "",
    "Verified SCPI commands:",
    "  CHANnel<m>:STATe ON|OFF",
    "  CHANnel<m>:SCALe <V/div>",
    "  CHANnel<m>:COUPling DCLimit|ACLimit|GND",
    "  TIMebase:SCALe <s/div>",
    "  TRIGger:A:SOURce CH1..CH4",
    "  TRIGger:A:LEVel1:VALue <volts>",
    "  TRIGger:A:EDGE:SLOPe POSitive|NEGative",
    "  MEASurement<m>:MAIN FREQuency|PEAK|RMS",
    "  MEASurement<m>:RESult:ACTual?",
], 14, WHITE)

bullet_box(s, 7, 2.3, 5.5, 4.8, [
    "Data flow:",
    "",
    "  Chat / Voice",
    "     ↓",
    "  Flask Backend (:5001)",
    "     ↓",
    "  Confirmation Gate (preview → approve)",
    "     ↓",
    "  RsInstrumentBackend",
    "     ↓",
    "  SCPI over USB TMC",
    "     ↓",
    "  R&S RTB24 Oscilloscope",
    "",
    "  Auto-discovery: list_resources('?*')",
    "  Fallback: simulator if hardware unavailable",
], 14, WHITE)

# ═══════════════════════════════════════════════
# SLIDE 9: Deliverables Checklist
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "TRACK 2 DELIVERABLES", 14, CYAN, True)
text_box(s, 0.8, 1.0, 11, 0.8, "All requirements met", 36, CYAN, True)

deliverables = [
    "✓  Live demo: 2 guided workflows (simple + safety-critical)",
    "✓  Assistant guides user step-by-step through procedures",
    "✓  Instrument configuration after explicit confirmation",
    "✓  Voice interaction (Gemini STT/TTS + GhanaNLP Twi)",
    "✓  Text chat interface (/chat.html)",
    "✓  Conversation transcripts (export JSON + server API)",
    "✓  Instructor view of logs and failure points",
    "✓  README with instrument/simulator connection guide",
    "✓  Code repository on GitHub",
]
bullet_box(s, 0.8, 2.2, 7, 4.5, deliverables, 18, WHITE)

stretch = [
    "Stretch goals achieved:",
    "",
    "✓  Voice integration with lip-synced avatar",
    "✓  Auto-configure instruments after confirmation",
    "✓  Twi language support (GhanaNLP)",
    "✓  AR component identification (Gemini Vision)",
    "✓  Anomaly detection with safety flags",
    "✓  Competency-based progress tracking",
]
bullet_box(s, 8.2, 2.2, 4.5, 4.5, stretch, 15, RGBColor(0xB0, 0xB8, 0xCC))

# ═══════════════════════════════════════════════
# SLIDE 10: Team
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 0.4, 10, 0.8, "TEAM", 14, CYAN, True)
text_box(s, 0.8, 1.2, 11, 1, "Four cores, one assistant", 36, WHITE, True)

teams = [
    ("Afari", "Stream I/O", "Mic/camera capture, STT/TTS pipeline,\navatar lip-sync, audio recorder"),
    ("Gregory", "RAG / Context", "15-document corpus, procedure/safety/fault\nknowledge base, sentence-transformer retrieval"),
    ("William", "Instrument API", "R&S RTB24 SCPI control, confirmation gate,\nRsInstrument integration, USB auto-discovery"),
    ("Prince", "Vision / Logging", "Component ID (Gemini Vision), anomaly detection,\nsession logging, competency tracking, progress system"),
]
for i, (name, role, desc) in enumerate(teams):
    x = 0.8 + (i % 2) * 6.2
    y = 2.8 + (i // 2) * 2.2
    text_box(s, x, y, 5.5, 0.5, name, 24, CYAN, True)
    text_box(s, x, y + 0.5, 5.5, 0.4, role, 16, WHITE, True)
    text_box(s, x, y + 0.9, 5.5, 0.8, desc, 13, DIM)

# ═══════════════════════════════════════════════
# SLIDE 11: Thank You
# ═══════════════════════════════════════════════
s = add_slide()
add_bar(s, 0, 0, 13.333, 0.06, CYAN)
text_box(s, 0.8, 2.0, 11, 1.5, "Thank You", 60, CYAN, True, PP_ALIGN.CENTER)
text_box(s, 0.8, 3.8, 11, 1, "Rhoda — Virtual Lab Assistant", 24, WHITE, False, PP_ALIGN.CENTER)
text_box(s, 0.8, 4.6, 11, 0.5, "github.com/PhidLarkson/rohde-schwarz-assistant", 16, DIM, False, PP_ALIGN.CENTER)
text_box(s, 0.8, 5.5, 11, 1, "Demo:  /chat.html  ·  XR Avatar  ·  /instructor.html\nBackend:  python server/app.py  |  SCOPE_RESOURCE=USB for real hardware", 14, DIM, False, PP_ALIGN.CENTER)

# Save
out = "docs/Rhoda_Hackathon_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
